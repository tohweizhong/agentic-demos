import os
import random
from docx import Document
from pptx import Presentation
from openpyxl import Workbook
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet

# Setup target folder
MOCK_DIR = os.path.dirname(os.path.abspath(__file__))

# Curated bank context dictionaries to generate high-fidelity, realistic bank files
BANK_TOPICS = {
    "risk": [
        {"title": "Credit Risk Assessment Framework", "details": "Evaluates counterparty exposure limits and default probabilities under high inflation scenarios."},
        {"title": "Liquidity Coverage Ratio LCR Audit", "details": "Reviews high-quality liquid assets against projected 30-day net cash outflows."},
        {"title": "Market Stress Testing Model", "details": "Models potential portfolio losses under simulated interest rate shocks of +200bps and +400bps."},
        {"title": "Operational Risk Incidents Log", "details": "Details transactional system outages, payment processing delays, and root cause remediation steps."},
        {"title": "Model Risk Management Policy", "details": "Outlines validation standards for automated valuation models and algorithmic trading bots."},
        {"title": "Interest Rate Risk in the Banking Book IRRBB Guide", "details": "Outlines the vulnerability of earnings and capital structure to adverse interest rate movements."},
        {"title": "Counterparty Credit Risk Management Policy", "details": "Manages exposures to OTC derivatives, securities financing, and clearing transactions."},
        {"title": "Climate Change Financial Risk Assessment Guidelines", "details": "Identifies transition and physical climate hazards on corporate loan assets."},
        {"title": "Sovereign Debt Exposure Limits Policy", "details": "Sets country-specific concentration caps for foreign national debt holdings."},
        {"title": "Fraud Prevention and Detection Analytics Model", "details": "Details machine learning pattern analysis for card transactions and account takeovers."}
    ],
    "compliance": [
        {"title": "Anti Money Laundering AML Standard Operating Procedures", "details": "Specifies Customer Due Diligence triggers, beneficial ownership verification, and suspicious transaction reporting rules."},
        {"title": "Sanctions Screening Guidelines", "details": "Ensures strict compliance with OFAC, EU, and local monetary authority restricted entity lists."},
        {"title": "Know Your Customer KYC Remediation Plan", "details": "Details the high-risk customer account review timeline and documentation requirements."},
        {"title": "Insider Trading Prevention Code", "details": "Defines trading blackout windows, pre-clearance procedures for employees, and conflict of interest disclosures."},
        {"title": "Personal Data Protection Policy PDPA", "details": "Mandates security safeguards for customer NRIC, credit cards, and account balances."},
        {"title": "Whistleblowing Investigation Procedures", "details": "Guarantees anonymous escalation channels and protects reporters from workplace retaliation."},
        {"title": "Anti Bribery and Corruption Compliance Manual", "details": "Governs interactions with foreign public officials, state-owned enterprises, and corporate gifts."},
        {"title": "Regulatory Reporting Submission Framework", "details": "Coordinates timely reporting to central banks, SEC, and tax authorities."},
        {"title": "Transaction Monitoring System Thresholds", "details": "Defines rule-based alerts for structuring, high-frequency transfers, and shell company activity."},
        {"title": "External Audit Engagement Guidelines", "details": "Establishes criteria for partnering with independent accounting firms and sharing files."}
    ],
    "wealth": [
        {"title": "High Net Worth Individual HNWI Portfolio Strategy", "details": "Recommends asset allocation adjustments favoring alternative investments, private equity, and green bonds."},
        {"title": "Private Banking Discretionary Mandates Guidelines", "details": "Details management rules for active balanced portfolios and tactical currency overlays."},
        {"title": "Sustainable Investing and ESG Criteria Overview", "details": "Defines screening standards for climate-resilient infrastructure projects and zero-carbon equities."},
        {"title": "Global Macroeconomic Outlook Q4 2026", "details": "Analyzes GDP growth rates, central bank policy pivots, and currency volatility impacts."},
        {"title": "Estate Planning and Trust Structures Guide", "details": "Explains generation-skipping trusts, offshore family offices, and tax minimization strategies."},
        {"title": "Family Office Tax Minimization Strategy", "details": "Evaluates international asset transfer legal vehicles to reduce estate tax liabilities."},
        {"title": "Structured Products and Derivatives Sales Handbook", "details": "Provides mandatory compliance criteria for presenting yield-enhancement notes to investors."},
        {"title": "Venture Capital Asset Class Allocation Guide", "details": "Establishes guidelines for financing early-stage fintech, AI startups, and biotech companies."},
        {"title": "Fixed Income Tactical Management Strategy", "details": "Directs corporate bond duration shifts to mitigate rate inflation impacts."},
        {"title": "Crypto Asset Advisory Framework for Wealth Management", "details": "Details limitations and risk warnings when recommending stablecoins and digital currencies."}
    ],
    "operations": [
        {"title": "Instant Payment Processing Gateway SOP", "details": "Details technical integration specifications, API endpoints, and transaction throughput optimization metrics."},
        {"title": "Core Banking Cloud Migration Plan", "details": "Outlines the multi-phase transition of ledger databases to secure hybrid cloud environments."},
        {"title": "SWIFT ISO 20022 Integration Standard", "details": "Specifies XML schema requirements for international cross-border payment messages."},
        {"title": "Business Continuity and Disaster Recovery Plan", "details": "Details backup server locations, failover hot-sites, and crisis communication team roles."},
        {"title": "Merchant Acquisition Onboarding Guide", "details": "Explains POS terminal configuration, settlement cycles, and merchant fee structures."},
        {"title": "Clearing and Settlement Operations Procedures", "details": "Establishes standard cutoff windows for ACH, Fedwire, and international payments."},
        {"title": "Cash Vault Replenishment and Security SOP", "details": "Outlines physical safety protocols, CIT schedules, and cash holding limit alerts."},
        {"title": "Nostro and Vostro Account Reconciliation Guide", "details": "Coordinates daily ledger matching with foreign correspondent banking partners."},
        {"title": "Credit Card Issuance and Personalization Workflow", "details": "Details chip embedding, security code activation, and mailing logistics."},
        {"title": "ATM Network Maintenance and Security Standard", "details": "Specifies physical inspection cycles and skimming prevention firmware updates."}
    ],
    "hr": [
        {"title": "Global Hybrid Work and Telecommuting Guidelines", "details": "Defines requirements for remote equipment, core working hours, and workspace ergonomics."},
        {"title": "Employee Offboarding and Security Exit Procedures", "details": "Coordinates access revocation, corporate asset return, and exit interview schedules."},
        {"title": "Parental Leave Benefits and Reintegration Policy", "details": "Outlines paid leave allowances for primary carers and return-to-work transition programs."},
        {"title": "Workplace Anti Harassment and Bullying Policy", "details": "Provides formal complaint channels and investigation protocols for employee misconduct."},
        {"title": "Global Mobility and Relocation Handbook", "details": "Details visa sponsorship, housing allowances, and tax assistance for international transfers."},
        {"title": "Employee Referral and Talent Acquisition Bonus Scheme", "details": "Outlines candidate eligibility, payment cycles, and referral rewards for core roles."},
        {"title": "Employee Mental Health and Wellness Program", "details": "Provides information on corporate counseling services, stress support, and gym subsidies."},
        {"title": "Workplace Health and Safety Standards", "details": "Mandates fire drills, first aid training, incident logging, and building safety reviews."},
        {"title": "Grievance Handling and Resolution Procedures", "details": "Guides employees on filing official complaints and outlines HR investigation processes."},
        {"title": "Retiree Benefit Continuation and Healthcare Policy", "details": "Specifies pension allocations and medical coverage provisions for long-service employees."}
    ],
    "procurement": [
        {"title": "Vendor Code of Conduct and Ethics Standard", "details": "Enforces human rights, anti-corruption, and eco-friendly requirements on all suppliers."},
        {"title": "RFP Technical and Commercial Evaluation Criteria", "details": "Provides detailed weighting matrices for vendor submissions, pricing, and architecture."},
        {"title": "Outsourcing Service Level Agreement Templates", "details": "Standardizes uptime penalties, performance reviews, and liability limitations."},
        {"title": "Sustainable Supply Chain Procurement Policy", "details": "Mandates green criteria for vendor selection, reducing plastic and energy usage."},
        {"title": "Vendor Onboarding and Due Diligence Checklist", "details": "Details bank details verification, cyber security assessments, and legal validation steps."},
        {"title": "Purchase Order Requisition and Approval Matrix", "details": "Specifies authorization thresholds, double-signature rules, and finance approval steps."},
        {"title": "Hardware Provisioning and Asset Lifecycles", "details": "Defines server procurement, developer laptop replacement cycles, and secure shredding."},
        {"title": "Enterprise Software License Renewal Framework", "details": "Establishes timelines for reviewing SaaS utility, contract negotiations, and exit plans."},
        {"title": "Vendor Penalty and Contract Termination Protocol", "details": "Governs breach of service response steps, default terms, and exit strategies."},
        {"title": "Independent Consultant Hiring and Billing Guidelines", "details": "Standardizes hourly caps, timesheet verification, and professional liability coverage."}
    ],
    "learning_development": [
        {"title": "Annual Mandatory Regulatory Compliance Training Policy", "details": "Outlines course completions for AML, insider trading, and cybersecurity for all staff."},
        {"title": "Executive Leadership and Management Development Program", "details": "Details training schedules for new directors, including communication and change management."},
        {"title": "Technical Engineering Upskilling and Cloud Pathway", "details": "Lists approved training for Kubernetes, AWS, and AI model engineering certifications."},
        {"title": "Tuition Reimbursement and Continuing Education Guidelines", "details": "Specifies eligibility criteria, approved degrees, and multi-year service bonds."},
        {"title": "Graduate Analyst Onboarding Curriculum", "details": "Outlines the rotation schedule, banking fundamentals bootcamps, and senior mentoring."},
        {"title": "Mentorship and Professional Coaching Framework", "details": "Pairs rising leaders with senior executives to drive retention and diversity."},
        {"title": "Agile and Scrum Master Certification Program", "details": "Tracks product management training across scrum, sprints, and agile workflow optimization."},
        {"title": "Executive Presentation and Media Training Guide", "details": "Covers public speaking, panel interviews, and press release communications coaching."},
        {"title": "Language Proficiency and Cross Border Training", "details": "Assists staff in key markets to master professional business writing and local idioms."},
        {"title": "Customer Empathy and High Value Service Training", "details": "Instructs relationship managers on communication, conflict resolution, and active listening."}
    ],
    "expense": [
        {"title": "Global Travel and Accommodation Expense Policy", "details": "Specifies flight booking tiers, hotel caps, and per-diem allowances by city grade."},
        {"title": "Client Entertainment and Corporate Hospitality Thresholds", "details": "Sets spending limits, business justification rules, and mandatory guest logs."},
        {"title": "Corporate Credit Card Issuance and Audit Rules", "details": "Defines candidate selection, monthly credit limits, and receipt submission deadlines."},
        {"title": "Capital Expenditure CAPEX Project Authorization Matrix", "details": "Details cost-benefit templates, ROI metrics, and CFO sign-off limits for new assets."},
        {"title": "Accounts Payable Invoice Processing and Matching SOP", "details": "Governs three-way matching of purchase orders, receiving logs, and supplier invoices."},
        {"title": "Departmental Budget Reallocation and Variance Guidelines", "details": "Instructs budget owners on shifting funds and reporting over-budget operations."},
        {"title": "Corporate Income Tax Filing and Compliance SOP", "details": "Coordinates local and foreign tax deductions, transfer pricing, and annual reporting."},
        {"title": "Treasury Liquidity Desk Operational Limits", "details": "Limits currency conversion ranges, overnight lending capacity, and short-term assets."},
        {"title": "Petty Cash Fund Management and Audit Standards", "details": "Maintains secure cash box thresholds, reconciliation logs, and spot audit cycles."},
        {"title": "Asset Depreciation and IT Write Off Policy", "details": "Coordinates with accounting to mark down obsolete servers, printers, and laptops."}
    ],
    "talent": [
        {"title": "Annual Performance Appraisal and Feedback Framework", "details": "Defines key evaluation benchmarks, rating structures, and employee self-assessments."},
        {"title": "Promotion Nomination and Technical Committee Review Standards", "details": "Specifies eligibility matrices, case submission packets, and interview rounds."},
        {"title": "Executive Succession Planning and Talent Review Protocol", "details": "Tracks critical roles, lists backups, and ensures immediate operational readiness."},
        {"title": "Multi Rater 360 Degree Feedback Process", "details": "Coordinates feedback from reports, peers, and directors to evaluate lead performance."},
        {"title": "High Potential HiPo Development Track Guidelines", "details": "Allocates funding for custom assignments, coaching, and fast-track promotion."},
        {"title": "Performance Improvement Plan PIP Guide for Managers", "details": "Standardizes timeline rules, goal definitions, and performance check schedules."},
        {"title": "Job Grading and Career Banding Criteria", "details": "Aligns titles from Analyst to Managing Director with competitive pay benchmarks."},
        {"title": "Lateral Transfer and Internal Mobility Framework", "details": "Allows staff to apply for roles in other divisions without management block."},
        {"title": "Equity Compensation and Stock Option Incentive Scheme", "details": "Controls vesting periods, executive share allocations, and buyout conditions."},
        {"title": "Engineering Career Path and Technical Contributor Tracks", "details": "Defines promotions for staff who want to focus on coding rather than management."}
    ],
    "it_security": [
        {"title": "User Access Password Policy and MFA Standard", "details": "Enforces complex keys, mandatory rotating cycles, and multi-factor authentication."},
        {"title": "Data Loss Prevention DLP Monitoring Rules", "details": "Blocks sharing of sensitive files, customer IDs, and code to external clouds."},
        {"title": "Bring Your Own Device BYOD Security Guidelines", "details": "Requires MDM installation, local storage encryption, and remote wipe rights."},
        {"title": "Phishing Simulation Drills and Reporting Requirements", "details": "Specifies email drill schedules, tracking metrics, and training defaults for clicks."},
        {"title": "Cybersecurity Incident Response and Escalation Plan", "details": "Outlines threat containment steps, forensic capture, and customer notice SLA."},
        {"title": "Server OS Patching and Vulnerability SLA", "details": "Enforces critical patching timeline (14 days) and regular code scans."},
        {"title": "Network Firewall Access Request and Rule Reviews", "details": "Governs approval processes for open ports and annual rules cleanup."},
        {"title": "Identity Access Management IAM Verification Cycle", "details": "Mandates access audits every six months to prune inactive or legacy roles."},
        {"title": "USB and Removable Storage Security Restriction", "details": "Disables write rights to standard USB flash drives on corporate machines."},
        {"title": "Helpdesk Escalation and System Incident Log Standard", "details": "Governs ticketing priorities from Level 1 troubleshooting to Major Incident teams."}
    ]
}

DEPARTMENTS = ["Retail Banking", "Wealth Management", "Corporate Investment Bank", "Risk Management", "Global Compliance", "Technology & Operations"]

def generate_docx(filepath, topic):
    doc = Document()
    doc.add_heading(topic["title"], level=1)
    doc.add_paragraph(f"Classification: Internal Only\nDepartment: {random.choice(DEPARTMENTS)}\nDate: 2026-06-01\n")
    
    p = doc.add_paragraph(
        f"This document establishes the official guidelines and framework regarding {topic['title'].lower()}. "
        f"It is designed to ensure all operational divisions remain strictly compliant with regional regulatory guidelines "
        f"and corporate governance policies.\n\n"
        f"Key Actionable Items & Guidelines:\n"
    )
    doc.add_paragraph(f"1. Implementation details: {topic['details']}", style='List Bullet')
    doc.add_paragraph("2. Regular audits must be conducted bi-annually to review progress and ensure absolute alignment.", style='List Bullet')
    doc.add_paragraph("3. Any operational exception requires written authorization from the Senior Managing Director.", style='List Bullet')
    doc.add_paragraph("4. Under Microsoft Purview rules, this document is labeled as General/All Employees.", style='List Bullet')
    
    doc.save(filepath)

def generate_pptx(filepath, topic):
    prs = Presentation()
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = topic["title"]
    subtitle.text = f"Strategic Overview & Framework\nDate: 2026-06-01 | {random.choice(DEPARTMENTS)}"
    
    # Slide 2: Details Slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Core Objectives & Scope"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = f"Operational mandate for {topic['title'].lower()}:"
    
    p1 = tf.add_paragraph()
    p1.text = f"- Focus: {topic['details']}"
    
    p2 = tf.add_paragraph()
    p2.text = "- Objective: Mitigate operational issues and optimize transaction efficiency."
    
    p3 = tf.add_paragraph()
    p3.text = "- Review Cycle: Monitored by the Executive Board quarterly."
    
    prs.save(filepath)

def generate_xlsx(filepath, topic):
    wb = Workbook()
    ws = wb.active
    ws.title = "Summary Metrics"
    
    # Headers
    ws.cell(row=1, column=1, value="Factual Metric Key")
    ws.cell(row=1, column=2, value="Value / Allocation")
    ws.cell(row=1, column=3, value="Audit Status")
    
    # Content rows
    ws.cell(row=2, column=1, value="Document Title Reference")
    ws.cell(row=2, column=2, value=topic["title"])
    ws.cell(row=2, column=3, value="Passed")
    
    ws.cell(row=3, column=1, value="Primary Detail Focus")
    ws.cell(row=3, column=2, value=topic["details"])
    ws.cell(row=3, column=3, value="Verified")
    
    ws.cell(row=4, column=1, value="Simulated Stress Coefficient")
    ws.cell(row=4, column=2, value=f"{round(random.uniform(1.1, 3.5), 2)}")
    ws.cell(row=4, column=3, value="Passed Stress")
    
    ws.cell(row=5, column=1, value="Target Annual Reserve Pool ($M)")
    ws.cell(row=5, column=2, value=random.randint(50, 750))
    ws.cell(row=5, column=3, value="Approved")
    
    wb.save(filepath)

def generate_pdf(filepath, topic):
    doc = SimpleDocTemplate(filepath, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    
    title_style = styles['Heading1']
    body_style = styles['BodyText']
    
    story.append(Paragraph(topic["title"], title_style))
    story.append(Spacer(1, 15))
    story.append(Paragraph(f"<b>Date:</b> 2026-06-01<br/><b>Department:</b> {random.choice(DEPARTMENTS)}", body_style))
    story.append(Spacer(1, 20))
    
    intro_text = (
        f"This corporate standard outlines the operational parameters for {topic['title'].lower()}. "
        f"Our central bank guidelines dictate rigorous auditing, strict model validation, and proactive risk assessment policies. "
        f"The objective is to maintain systemic resilience and support financial markets safely."
    )
    story.append(Paragraph(intro_text, body_style))
    story.append(Spacer(1, 15))
    
    details_text = f"<b>Core Implementation Strategy:</b> {topic['details']}"
    story.append(Paragraph(details_text, body_style))
    story.append(Spacer(1, 15))
    
    conclusion = "All staff members are instructed to review these materials and complete the required compliance verification online."
    story.append(Paragraph(conclusion, body_style))
    
    doc.build(story)

def clean_old_mock_files():
    """Deletes any existing mock files to clean the directory."""
    for f in os.listdir(MOCK_DIR):
        if f.split(".")[-1] in ["docx", "pptx", "xlsx", "pdf"]:
            try:
                os.remove(os.path.join(MOCK_DIR, f))
            except Exception:
                pass

def main():
    print("Cleaning up old generic mock files...")
    clean_old_mock_files()
    
    print("Kicking off corporate mock file generation script with meaningful naming...")
    
    formats = ["docx", "pptx", "xlsx", "pdf"]
    total_files_needed = 100
    
    # Gather all topics flatly
    flat_topics = []
    for cat, topics in BANK_TOPICS.items():
        for t in topics:
            flat_topics.append(t)
            
    file_count = 0
    
    while file_count < total_files_needed:
        topic = flat_topics[file_count % len(flat_topics)]
        version = file_count // len(flat_topics) + 1
        
        topic_copy = {
            "title": f"{topic['title']} - Version {version}",
            "details": topic["details"]
        }
        
        # Make a safe, descriptive filename
        safe_title = topic["title"].replace(" ", "_").replace("(", "").replace(")", "").replace("-", "_").replace("/", "_")
        file_ext = formats[file_count % len(formats)]
        filename = f"{safe_title}_v{version}.{file_ext}"
        filepath = os.path.join(MOCK_DIR, filename)
        
        if file_ext == "docx":
            generate_docx(filepath, topic_copy)
        elif file_ext == "pptx":
            generate_pptx(filepath, topic_copy)
        elif file_ext == "xlsx":
            generate_xlsx(filepath, topic_copy)
        elif file_ext == "pdf":
            generate_pdf(filepath, topic_copy)
            
        file_count += 1
        if file_count % 10 == 0:
            print(f"Generated {file_count} mock files...")
            
    print(f"Successfully created {total_files_needed} corporate-themed files with meaningful names in {MOCK_DIR}!")

if __name__ == "__main__":
    main()
